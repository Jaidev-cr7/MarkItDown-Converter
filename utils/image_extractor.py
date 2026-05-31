"""
image_extractor.py – Extract embedded images from documents for OCR.

Returns a list of ExtractedImage(bytes, position_label) objects.
Each caller (PDF, DOCX, PPTX, HTML) is isolated and degrades gracefully
if its dependency is absent.

Supported formats
-----------------
- PDF   → PyMuPDF (fitz)
- DOCX  → python-docx
- PPTX  → python-pptx  (already available via markitdown[all])
- HTML  → BeautifulSoup4 (already available via markitdown[all])
- Image → pass-through (the file itself is the image)

No Tesseract dependency.  All OCR is handled by ocr_handler.py.
"""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Sequence

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Availability flags – graceful degradation
# ---------------------------------------------------------------------------

try:
    import fitz  # PyMuPDF  # type: ignore
    _FITZ_OK = True
except ImportError:
    _FITZ_OK = False

try:
    import docx  # python-docx  # type: ignore
    _DOCX_OK = True
except ImportError:
    _DOCX_OK = False

try:
    import pptx  # python-pptx  # type: ignore
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

try:
    from bs4 import BeautifulSoup  # type: ignore
    import urllib.request
    _BS4_OK = True
except ImportError:
    _BS4_OK = False

try:
    from PIL import Image as _PILImage  # type: ignore
    _PIL_OK = True
except ImportError:
    _PIL_OK = False


# ---------------------------------------------------------------------------
# Data class
# ---------------------------------------------------------------------------

@dataclass
class ExtractedImage:
    """A single image extracted from a document."""

    data: bytes                  # Raw image bytes (PNG or JPEG)
    mime: str = "image/png"      # MIME type hint
    position_label: str = ""     # Human-readable location hint, e.g. "Page 2" or "Slide 3"
    source_name: str = ""        # Original embedded name / xref hint, if known

    # Computed lazily
    _sha256: str = field(default="", init=False, repr=False)

    @property
    def sha256(self) -> str:
        if not self._sha256:
            import hashlib
            self._sha256 = hashlib.sha256(self.data).hexdigest()
        return self._sha256


# ---------------------------------------------------------------------------
# Public dispatcher
# ---------------------------------------------------------------------------

IMAGE_EXTENSIONS: frozenset[str] = frozenset(
    {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tiff", ".tif"}
)


def extract_images(path: Path) -> list[ExtractedImage]:
    """Return all embedded images found in *path*.

    Dispatches to the appropriate format handler.  Returns an empty list
    if the format is not supported or the dependency is unavailable.
    """
    suffix = path.suffix.lower()
    handler: Callable[[Path], list[ExtractedImage]] | None = None

    if suffix == ".pdf":
        handler = _extract_pdf
    elif suffix in {".docx", ".doc"}:
        handler = _extract_docx
    elif suffix in {".pptx", ".ppt"}:
        handler = _extract_pptx
    elif suffix in {".html", ".htm"}:
        handler = _extract_html
    elif suffix in IMAGE_EXTENSIONS:
        handler = _passthrough_image
    # xlsx, csv, txt, json, xml, epub, audio → no embedded images to extract

    if handler is None:
        return []

    try:
        return handler(path)
    except Exception:
        logger.exception("image_extractor: failed to extract images from %s", path.name)
        return []


# ---------------------------------------------------------------------------
# PDF extractor  (PyMuPDF)
# ---------------------------------------------------------------------------

# Minimum image dimensions to bother OCR-ing (avoids tiny icons/bullets)
_MIN_IMG_WIDTH  = 60
_MIN_IMG_HEIGHT = 60
# Minimum byte size – very small blobs are likely vector shapes rendered as 1px
_MIN_IMG_BYTES  = 500


def _extract_pdf(path: Path) -> list[ExtractedImage]:
    if not _FITZ_OK:
        logger.debug("PyMuPDF not available – skipping PDF image extraction")
        return []

    images: list[ExtractedImage] = []
    doc = fitz.open(str(path))

    for page_idx, page in enumerate(doc, start=1):
        label = f"Page {page_idx}"
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
            except Exception:
                continue

            img_bytes: bytes = base_image["image"]
            w: int = base_image.get("width", 0)
            h: int = base_image.get("height", 0)
            ext: str = base_image.get("ext", "png")

            if w < _MIN_IMG_WIDTH or h < _MIN_IMG_HEIGHT:
                continue
            if len(img_bytes) < _MIN_IMG_BYTES:
                continue

            mime = "image/jpeg" if ext in ("jpg", "jpeg") else "image/png"
            images.append(
                ExtractedImage(
                    data=img_bytes,
                    mime=mime,
                    position_label=label,
                    source_name=f"xref_{xref}",
                )
            )

    doc.close()
    return images


# ---------------------------------------------------------------------------
# DOCX extractor  (python-docx)
# ---------------------------------------------------------------------------

def _extract_docx(path: Path) -> list[ExtractedImage]:
    if not _DOCX_OK:
        logger.debug("python-docx not available – skipping DOCX image extraction")
        return []

    images: list[ExtractedImage] = []
    try:
        document = docx.Document(str(path))
    except Exception:
        logger.exception("Failed to open DOCX: %s", path.name)
        return []

    # Iterate paragraphs in order to preserve approximate position
    para_idx = 0
    image_idx = 0
    for para in document.paragraphs:
        para_idx += 1
        for run in para.runs:
            for element in run._element:
                # Inline image references: <a:blip r:embed="rId...">
                tag = getattr(element, "tag", "") or ""
                if "graphicData" in tag or "drawing" in tag or "pict" in tag:
                    # Drill into the relationship parts
                    for rel in document.part.rels.values():
                        if "image" in rel.reltype:
                            try:
                                img_bytes = rel.target_part.blob
                                if len(img_bytes) >= _MIN_IMG_BYTES:
                                    image_idx += 1
                                    _append_image(
                                        images,
                                        img_bytes,
                                        position_label=f"Paragraph {para_idx}",
                                        source_name=f"image_{image_idx}",
                                    )
                            except Exception:
                                pass
                    break  # Only one image per drawing element

    # Fallback: if the paragraph traversal missed images, pull all rels
    if not images:
        image_idx = 0
        seen: set[str] = set()
        for rel_id, rel in document.part.rels.items():
            if "image" in rel.reltype:
                try:
                    img_bytes = rel.target_part.blob
                    if len(img_bytes) >= _MIN_IMG_BYTES and rel_id not in seen:
                        seen.add(rel_id)
                        image_idx += 1
                        _append_image(
                            images,
                            img_bytes,
                            position_label=f"Embedded image {image_idx}",
                            source_name=rel_id,
                        )
                except Exception:
                    pass

    return images


# ---------------------------------------------------------------------------
# PPTX extractor  (python-pptx)
# ---------------------------------------------------------------------------

def _extract_pptx(path: Path) -> list[ExtractedImage]:
    if not _PPTX_OK:
        logger.debug("python-pptx not available – skipping PPTX image extraction")
        return []

    images: list[ExtractedImage] = []
    try:
        presentation = pptx.Presentation(str(path))
    except Exception:
        logger.exception("Failed to open PPTX: %s", path.name)
        return []

    from pptx.util import Pt  # noqa: F401 – just confirming pptx is usable
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

    for slide_idx, slide in enumerate(presentation.slides, start=1):
        label = f"Slide {slide_idx}"
        for shape in slide.shapes:
            # Direct picture shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    if len(img_bytes) >= _MIN_IMG_BYTES:
                        _append_image(
                            images,
                            img_bytes,
                            position_label=label,
                            source_name=getattr(shape, "name", ""),
                        )
                except Exception:
                    pass
            # Placeholder shapes that contain images
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                try:
                    if hasattr(shape, "image"):
                        img_bytes = shape.image.blob
                        if len(img_bytes) >= _MIN_IMG_BYTES:
                            _append_image(
                                images,
                                img_bytes,
                                position_label=label,
                                source_name=getattr(shape, "name", ""),
                            )
                except Exception:
                    pass

    return images


# ---------------------------------------------------------------------------
# HTML extractor  (BeautifulSoup4)
# ---------------------------------------------------------------------------

def _extract_html(path: Path) -> list[ExtractedImage]:
    if not _BS4_OK:
        logger.debug("BeautifulSoup4 not available – skipping HTML image extraction")
        return []

    html_text = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html_text, "html.parser")
    images: list[ExtractedImage] = []
    base_dir = path.parent

    for img_idx, tag in enumerate(soup.find_all("img"), start=1):
        src = tag.get("src", "").strip()
        if not src:
            continue

        img_bytes: bytes | None = None

        # Data-URI embedded images
        if src.startswith("data:image"):
            try:
                import base64
                _, encoded = src.split(",", 1)
                img_bytes = base64.b64decode(encoded)
            except Exception:
                pass
        # Local file references
        elif not src.startswith("http"):
            local_path = base_dir / src
            if local_path.exists():
                try:
                    img_bytes = local_path.read_bytes()
                except Exception:
                    pass
        # Remote URLs (best-effort, only if urllib available)
        else:
            try:
                req = urllib.request.Request(
                    src,
                    headers={"User-Agent": "Mozilla/5.0"},
                )
                with urllib.request.urlopen(req, timeout=5) as resp:
                    img_bytes = resp.read()
            except Exception:
                pass

        if img_bytes and len(img_bytes) >= _MIN_IMG_BYTES:
            _append_image(
                images,
                img_bytes,
                position_label=f"Image {img_idx}",
                source_name=src[:80],
            )

    return images


# ---------------------------------------------------------------------------
# Pass-through for standalone image files
# ---------------------------------------------------------------------------

def _passthrough_image(path: Path) -> list[ExtractedImage]:
    """The document IS the image — return its bytes as a single-element list."""
    try:
        img_bytes = path.read_bytes()
        if len(img_bytes) < _MIN_IMG_BYTES:
            return []
        return [
            ExtractedImage(
                data=img_bytes,
                mime="image/png",
                position_label="Full image",
                source_name=path.name,
            )
        ]
    except Exception:
        return []


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _append_image(
    out: list[ExtractedImage],
    img_bytes: bytes,
    position_label: str,
    source_name: str,
) -> None:
    """Normalise to PNG via Pillow (if available) and append to *out*."""
    if not _PIL_OK:
        # No normalisation possible; append raw bytes and hope for the best
        out.append(
            ExtractedImage(
                data=img_bytes,
                mime="image/png",
                position_label=position_label,
                source_name=source_name,
            )
        )
        return

    try:
        buf = io.BytesIO(img_bytes)
        img = _PILImage.open(buf)

        # Skip tiny images (icons, bullets, spacers)
        w, h = img.size
        if w < _MIN_IMG_WIDTH or h < _MIN_IMG_HEIGHT:
            return

        # Convert to RGB PNG for EasyOCR compatibility
        out_buf = io.BytesIO()
        img.convert("RGB").save(out_buf, format="PNG")
        out.append(
            ExtractedImage(
                data=out_buf.getvalue(),
                mime="image/png",
                position_label=position_label,
                source_name=source_name,
            )
        )
    except Exception:
        logger.debug("Could not normalise image '%s' via Pillow", source_name)
